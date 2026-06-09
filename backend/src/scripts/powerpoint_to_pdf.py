import fitz
import sys
import pptx

def powerpoint_to_pdf(input_path, output_path):
    try:
        prs = pptx.Presentation(input_path)
    except Exception as e:
        print(f"Error loading PowerPoint presentation: {e}", file=sys.stderr)
        sys.exit(1)
        
    html_content = "<html><body style='font-family: Arial, sans-serif; background-color: #f1f5f9; padding: 20px; color: #1e293b;'>"
    
    for idx, slide in enumerate(prs.slides):
        html_content += "<div style='background-color: white; border: 1px solid #cbd5e1; border-radius: 8px; padding: 30px; margin-bottom: 20px; box-shadow: 0 4px 6px rgba(0,0,0,0.05); min-height: 450px; page-break-after: always;'>"
        html_content += f"<div style='color: #94a3b8; font-size: 11px; margin-bottom: 15px; font-weight: bold; text-transform: uppercase; letter-spacing: 0.05em;'>Slide {idx + 1}</div>"
        
        slide_title = ""
        slide_paragraphs = []
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if not text:
                continue
            
            # Identify title placeholders
            is_title = False
            if shape.name.startswith("Title") or (hasattr(shape, "is_placeholder") and shape.is_placeholder and shape.placeholder_format.type in [1, 3]):
                is_title = True
                
            if is_title:
                slide_title = text
            else:
                for p in shape.text_frame.paragraphs:
                    p_text = p.text.strip()
                    if p_text:
                        slide_paragraphs.append((p.level, p_text))
                        
        if slide_title:
            html_content += f"<h1 style='color: #1e3a8a; font-size: 20px; margin-top: 0; margin-bottom: 20px; border-bottom: 1.5px solid #cbd5e1; padding-bottom: 8px;'>{slide_title}</h1>"
        else:
            html_content += "<div style='height: 20px;'></div>"
            
        if slide_paragraphs:
            for level, text in slide_paragraphs:
                indent = f"margin-left: {level * 24}px;" if level > 0 else ""
                prefix = "• " if level > 0 or len(slide_paragraphs) > 1 else ""
                font_size = "13px" if level == 0 else "12px"
                color = "#334155" if level == 0 else "#475569"
                html_content += f"<p style='font-size: {font_size}; color: {color}; margin: 6px 0; {indent}'>{prefix}{text}</p>"
        else:
            html_content += "<p style='color: #94a3b8; font-style: italic; font-size: 12px;'>Empty Slide</p>"
            
        html_content += "</div>"
        
    html_content += "</body></html>"
    
    try:
        fitz_doc = fitz.open(stream=html_content.encode('utf-8'), filetype='html')
        pdf_bytes = fitz_doc.convert_to_pdf()
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        fitz_doc.close()
    except Exception as e:
        print(f"Error generating PDF: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python powerpoint_to_pdf.py <input_path> <output_path>", file=sys.stderr)
        sys.exit(1)
    powerpoint_to_pdf(sys.argv[1], sys.argv[2])
