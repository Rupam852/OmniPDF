import fitz
import sys
import os

def pdf_to_pptx_gemini(input_path, output_path, gemini_key):
    try:
        import google.generativeai as genai
        import pptx
    except ImportError:
        print("Required python packages for Gemini conversion not found, using direct conversion.", file=sys.stderr)
        pdf_to_pptx_direct(input_path, output_path)
        return

    try:
        genai.configure(api_key=gemini_key)
        with open(input_path, "rb") as f:
            pdf_data = f.read()
            
        model = genai.GenerativeModel('gemini-2.5-flash')
        prompt = (
            "Analyze this PDF document and structure it as a slide-by-slide presentation outline. "
            "For each slide, output the slide title followed by bullet points. Format it clearly like:\n"
            "Slide 1 Title\n- point 1\n- point 2\n\nSlide 2 Title\n- point 1\n- point 2"
        )
        
        response = model.generate_content([
            {
                "mime_type": "application/pdf",
                "data": pdf_data
            },
            prompt
        ])
        text = response.text
    except Exception as e:
        print(f"Gemini API conversion failed: {e}. Falling back to direct conversion...", file=sys.stderr)
        pdf_to_pptx_direct(input_path, output_path)
        return

    try:
        prs = pptx.Presentation()
        # Set slide size to widescreen 16:9
        prs.slide_width = pptx.util.Inches(13.333)
        prs.slide_height = pptx.util.Inches(7.5)
        
        slides_raw = text.split("\n\n")
        for slide_data in slides_raw:
            lines = [l.strip() for l in slide_data.split("\n") if l.strip()]
            if not lines:
                continue
                
            title = lines[0]
            if title.lower().startswith("slide "):
                # Trim e.g., "Slide 1: Title"
                title_parts = title.split(":", 1)
                if len(title_parts) > 1:
                    title = title_parts[1].strip()
                else:
                    # check space
                    parts_space = title.split(" ", 2)
                    if len(parts_space) > 2:
                        title = parts_space[2].strip()
            
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            
            tf = slide.placeholders[1].text_frame
            tf.clear()
            
            for point in lines[1:]:
                if point.startswith("- ") or point.startswith("* "):
                    p = tf.add_paragraph()
                    p.text = point[2:]
                    p.level = 0
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 1
                    
        prs.save(output_path)
    except Exception as e:
        print(f"Error saving PPTX from Gemini output: {e}. Falling back to direct conversion...", file=sys.stderr)
        pdf_to_pptx_direct(input_path, output_path)

def pdf_to_pptx_direct(input_path, output_path):
    try:
        import pptx
        prs = pptx.Presentation()
        prs.slide_width = pptx.util.Inches(10)
        prs.slide_height = pptx.util.Inches(7.5)
        
        doc = fitz.open(input_path)
        for idx, page in enumerate(doc):
            text = page.get_text().strip()
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            if lines:
                slide.shapes.title.text = lines[0]
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for line in lines[1:10]: # limit to 10 lines to avoid overflow
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
            else:
                slide.shapes.title.text = f"Slide {idx + 1}"
                
        prs.save(output_path)
        doc.close()
    except Exception as e:
        print(f"Direct PDF to PPTX conversion failed: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python pdf_to_powerpoint.py <input_path> <output_path> [gemini_api_key]", file=sys.stderr)
        sys.exit(1)
        
    in_p = sys.argv[1]
    out_p = sys.argv[2]
    key = sys.argv[3] if len(sys.argv) > 3 and sys.argv[3].strip() and sys.argv[3] != "undefined" else ""
    
    if key:
        pdf_to_pptx_gemini(in_p, out_p, key)
    else:
        pdf_to_pptx_direct(in_p, out_p)
